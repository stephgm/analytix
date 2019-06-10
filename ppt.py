# -*- coding: utf-8 -*-
"""
Created on Wed Apr 17 17:39:41 2019

@author: Jordan
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from lxml import etree

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
chart_data.add_series('Q3 Salggggggggges', (20.4, 26.3, 14.2))
chart_data.add_series('Q4 SALES',(1,2,3))
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart
chart.has_legend = True
# This call doesn't fail, but also doesn't work
#chart.has_data_table = True

#XML for data table
dTable= """
        <c:dTable xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
            <c:showHorzBorder val="1"/>
            <c:showVertBorder val="1"/>
            <c:showOutline val="1"/>
            <c:showKeys val="1"/>
            <c:spPr>
                <a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:solidFill>
                        <a:schemeClr val="tx1"/>
                    </a:solidFill>
                </a:ln>
            </c:spPr>
        </c:dTable>
        """
#XML for Chart Border
ChartBorder = """
        <c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
            <a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />
            <a:ln cap="sq" w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" >
                <a:solidFill>
                    <a:schemeClr val="tx1">
                        <a:lumMod val="95000"/>
                        <a:lumOff val="5000"/>
                    </a:schemeClr>
                </a:solidFill>
            <a:bevel/>
            </a:ln>
        </c:spPr>
"""
#XML for Legend Border

LegendBorder = """
        <c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"> 
           <a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" >   
                <a:schemeClr val="bg1"/>   
           </a:solidFill> 

           <a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"> 
               <a:solidFill> 
                   <a:schemeClr val="tx1"/> 
               </a:solidFill> 
           </a:ln> 
        </c:spPr> 
"""

#This definitely works, but there is no way to get the damn axis extent as of yet.
#And its also impractical due to Legend names lengths being unknown, therefore the 
# width cannot be automatically determined and also it would require hamjamming... just drag the damn legend.

#Coordinates 1,1 correspond to bottom right of the Chart, however, this is not the bottom right of Axis.
# the subtracted numbers are the width and height of the legend respectively... currently no way to calc.
legendx = 1.-0.21850138524351123
legendy = 1.-0.3426567512394284
LegendPos = """
        
        <c:layout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"> 
            <c:manualLayout> 
                <c:xMode val="edge"/> 
                <c:yMode val="edge"/> 
                <c:x x="%s"/> 
                <c:y y="%s"/> 
                <c:w val="0.21850138524351123"/> 
                <c:h val="0.3426567512394284"/> 
            </c:manualLayout> 
        </c:layout>
""" % (legendx,legendy)


chart.plots._plotArea.append(etree.fromstring(dTable))
chart.plots._plotArea.append(etree.fromstring(ChartBorder))
chart.legend._element.append(etree.fromstring(LegendBorder))
#chart.legend._element.append(etree.fromstring(LegendPos))
prs.save('chart-01.pptx')